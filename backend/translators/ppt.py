from .base import Translator
from pptx import Presentation


class PPTTranslator(Translator):
    def translate_ppt(self, file_path: str, output_path: str):
        ppt = Presentation(file_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    original_text = shape.text
                    translated_text = self.translate_text(original_text)
                    shape.text = translated_text

        ppt.save(output_path)
